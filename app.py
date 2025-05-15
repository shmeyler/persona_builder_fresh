
def extract_insights_from_text_block(text_block):
    import re
    insights = []
    current = {}

    for line in text_block.splitlines():
        line = line.strip()

        if not line or re.match(r"^[0-9]{2}\.[0-9]{2}\.[0-9]{4}", line):  # Skip date stamps
            continue

        if re.match(r"^[A-Z ]{3,}$", line):  # Category like "TICKETING SERVICES USED"
            current["Category"] = line.title()
            continue

        if line.endswith("%"):
            try:
                comp = float(line.strip('%'))
                current["Composition"] = comp
            except:
                continue
        elif re.match(r"^\d{2,3}$", line):
            current["Index"] = int(line)
        elif len(line.split()) <= 6:
            current["Insight"] = line

        if len(current) >= 3:
            insights.append(current)
            current = {}

    return insights


import streamlit as st
import os
import io
import logging
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.cloud import vision
from PIL import Image
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import json
import csv

# === SECTION 1: Resonate UI Mapping and Prepopulation ===
resonate_taxonomy_map = {
    "Values and Beliefs": {
        "category": "Personal Values",
        "attributes": [
            "Stimulation", "Tolerance", "Achievement", "Equality", "Independence", "Tradition", "Security"
        ]
    },
    "Motivations": {
        "category": "Psychological Drivers",
        "attributes": [
            "Life as an adventure", "Freedom", "Recognition", "Safety", "Curiosity", "Being in control"
        ]
    },
    "Pain Points": {
        "category": "Purchase Barriers",
        "attributes": [
            "Cost", "Trust", "Complexity", "Time commitment", "Lack of customization", "Confusing tech"
        ]
    },
    "Media Habits": {
        "category": "Media Consumption",
        "attributes": [
            "Streaming", "Social Media", "TV", "Podcasts", "Mobile Apps", "News Sites"
        ]
    },
    "Goals and Aspirations": {
        "category": "Life Goals",
        "attributes": [
            "Career advancement", "Work-life balance", "Helping others", "Wealth", "Adventure", "Stability"
        ]
    }
}

def extract_resonate_defaults(text_blocks, taxonomy_map):
    joined_text = "\n".join(text_blocks).lower()
    defaults = {}
    for field, meta in taxonomy_map.items():
        matches = [attr for attr in meta["attributes"] if attr.lower() in joined_text]
        if matches:
            defaults[field] = matches
    return defaults

def build_persona_form_ui(mapping, defaults=None):
    st.header("Refine Persona Attributes")
    results = {}
    for field, meta in mapping.items():
        with st.expander(f"{field} ({meta['category']})"):
            suggestions = meta["attributes"]
            default_vals = defaults.get(field, []) if defaults else []
            selected = st.multiselect(f"Select from Resonate Attributes for '{field}'", suggestions, default=default_vals)

            manual_input = st.text_input(f"Optional: Add custom entries for '{field}' (comma-separated)", key=field)
            manual_entries = [x.strip() for x in manual_input.split(",") if x.strip()]
            matched = [entry for entry in manual_entries if entry in suggestions]
            unmatched = [entry for entry in manual_entries if entry not in suggestions]

            results[field] = {
                "selected": selected,
                "manual": manual_entries,
                "matched": matched,
                "unmatched": unmatched
            }

            if unmatched:
                st.warning(f"Unmatched entries: {', '.join(unmatched)}")
    return results

# === SECTION 2: File Processing + Persona Builder ===
st.set_page_config(page_title="Persona Builder", layout="wide")
st.title("Persona Builder")

folder_id = st.text_input("Enter Google Drive Folder ID to scan:")

def load_gdrive_service():
    credentials = service_account.Credentials.from_service_account_info(
        st.secrets["gcp"],
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )
    return build("drive", "v3", credentials=credentials)

def load_vision_client():
    credentials = service_account.Credentials.from_service_account_info(st.secrets["gcp"])
    return vision.ImageAnnotatorClient(credentials=credentials)

def list_drive_files(service, folder_id):
    files = []
    query = f"'{folder_id}' in parents and trashed = false"
    page_token = None
    while True:
        response = service.files().list(
            q=query,
            spaces="drive",
            fields="nextPageToken, files(id, name, mimeType)",
            pageToken=page_token,
        ).execute()
        for file in response.get("files", []):
            if file["mimeType"] == "application/vnd.google-apps.folder":
                files += list_drive_files(service, file["id"])
            else:
                files.append(file)
        page_token = response.get("nextPageToken", None)
        if page_token is None:
            break
    return files

def download_file(service, file_id):
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        status, done = downloader.next_chunk()
    fh.seek(0)
    return fh

def extract_text_from_image(image_bytes, vision_client):
    image = vision.Image(content=image_bytes)
    response = vision_client.text_detection(image=image)
    if response.error.message:
        raise Exception(response.error.message)
    return response.text_annotations[0].description if response.text_annotations else ""

if folder_id:
    st.write(f"📁 Scanning folder `{folder_id}`...")
    try:
        drive_service = load_gdrive_service()
        vision_client = load_vision_client()
        all_files = list_drive_files(drive_service, folder_id)
        st.success(f"✅ Found {len(all_files)} files")

        all_text_blocks = []

        for file in all_files:
            name = file["name"]
            mime = file["mimeType"]
            st.write(f"🔍 Processing `{name}` ({mime})")
            try:
                if mime.startswith("image/"):
                    img_bytes = download_file(drive_service, file["id"]).read()
                    text = extract_text_from_image(img_bytes, vision_client)
                    st.text_area(f"📄 Extracted text from {name}:", text, height=200)
                    all_text_blocks.append(text)

                elif mime == "text/csv":
                    file_bytes = download_file(drive_service, file["id"])
                    df = pd.read_csv(file_bytes)
                    st.dataframe(df.head())

                elif mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    file_bytes = download_file(drive_service, file["id"])
                    df = pd.read_excel(file_bytes)
                    st.dataframe(df.head())

                elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file_bytes = download_file(drive_service, file["id"])
                    doc = Document(file_bytes)
                    text = "\n".join([para.text for para in doc.paragraphs])
                    st.text_area(f"📝 DOCX Contents: {name}", text, height=200)
                    all_text_blocks.append(text)

                elif mime == "application/pdf":
                    file_bytes = download_file(drive_service, file["id"])
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    st.text_area(f"📄 PDF Contents: {name}", text, height=200)
                    all_text_blocks.append(text)

                elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file_bytes = download_file(drive_service, file["id"])
                    prs = Presentation(file_bytes)
                    text_runs = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_runs.append(shape.text)
                    full_text = "\n\n".join(text_runs)
                    st.text_area(f"📽️ PPTX Contents: {name}", full_text, height=200)
                    all_text_blocks.append(full_text)

                else:
                    st.warning(f"⏭️ Unsupported MIME type: {mime}")

            except Exception as file_err:
                st.error(f"❌ Failed to process {name}: {file_err}")

        if all_text_blocks:
            st.divider()
            st.subheader("🧠 Now Refine the Persona with Resonate Mapping")

            defaults = extract_resonate_defaults(all_text_blocks, resonate_taxonomy_map)
            st.info("✨ Prepopulated suggestions based on document scan.")
            form_data = build_persona_form_ui(resonate_taxonomy_map, defaults)

            st.subheader("📋 Final Persona Summary")
            persona_summary = {}
            for field, values in form_data.items():
                combined = list(set(values["selected"] + values["matched"] + values["manual"]))
                persona_summary[field] = ", ".join(combined)
            st.table(persona_summary.items())

            json_str = json.dumps(persona_summary, indent=2)
            st.download_button("📥 Download as JSON", json_str, file_name="persona.json", mime="application/json")

            csv_buffer = io.StringIO()
            writer = csv.writer(csv_buffer)
            writer.writerow(["Field", "Value"])
            for k, v in persona_summary.items():
                writer.writerow([k, v])
            st.download_button("📥 Download as CSV", csv_buffer.getvalue(), file_name="persona.csv", mime="text/csv")

    except Exception as e:
        st.error(f"❌ Error initializing services or accessing folder: {e}")
